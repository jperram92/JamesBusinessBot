import logging
import os
from datetime import datetime
from typing import Dict, List

from docx import Document
from pptx import Presentation
from jinja2 import Environment, FileSystemLoader

logger = logging.getLogger(__name__)

class DocumentService:
    """Document generation service implementation."""
    
    def __init__(self, config: Dict):
        """Initialize the document service with configuration."""
        self.config = config
        self.template_dir = config['documents']['template_dir']
        self.output_dir = config['documents']['output_dir']
        
        # Initialize Jinja2 environment
        self.env = Environment(
            loader=FileSystemLoader(self.template_dir),
            autoescape=True
        )
        
        # Create output directory if it doesn't exist
        os.makedirs(self.output_dir, exist_ok=True)
        
        logger.info("Document service initialized")
    
    async def create_word_document(self, meeting_id: str, data: Dict) -> str:
        """Create a Word document from meeting data."""
        try:
            # Load template
            template = self.env.get_template('meeting_summary.docx')
            
            # Render template with data
            content = template.render(**data)
            
            # Create new document
            doc = Document()
            
            # Add content to document
            for paragraph in content.split('\n'):
                doc.add_paragraph(paragraph)
            
            # Save document
            filename = f"meeting_summary_{meeting_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
            filepath = os.path.join(self.output_dir, filename)
            doc.save(filepath)
            
            logger.info(f"Created Word document: {filepath}")
            return filepath
        except Exception as e:
            logger.error(f"Failed to create Word document: {str(e)}")
            raise
    
    async def create_powerpoint_presentation(self, meeting_id: str, data: Dict) -> str:
        """Create a PowerPoint presentation from meeting data."""
        try:
            # Load template
            template = self.env.get_template('meeting_presentation.pptx')
            
            # Render template with data
            content = template.render(**data)
            
            # Create new presentation
            prs = Presentation()
            
            # Add slides based on content
            slides_content = content.split('\n\n')
            for slide_content in slides_content:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                body = slide.shapes.placeholders[1]
                
                # Split content into title and body
                lines = slide_content.split('\n')
                title.text = lines[0]
                
                # Add body content
                tf = body.text_frame
                for line in lines[1:]:
                    p = tf.add_paragraph()
                    p.text = line
            
            # Save presentation
            filename = f"meeting_presentation_{meeting_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            filepath = os.path.join(self.output_dir, filename)
            prs.save(filepath)
            
            logger.info(f"Created PowerPoint presentation: {filepath}")
            return filepath
        except Exception as e:
            logger.error(f"Failed to create PowerPoint presentation: {str(e)}")
            raise
    
    def _format_action_items(self, action_items: List[Dict]) -> str:
        """Format action items for document inclusion."""
        formatted_items = []
        for item in action_items:
            formatted_item = f"- {item['description']}\n  Assigned to: {item['assignee']}\n  Due date: {item['due_date']}"
            formatted_items.append(formatted_item)
        return '\n'.join(formatted_items)
    
    def _format_key_points(self, key_points: List[str]) -> str:
        """Format key points for document inclusion."""
        return '\n'.join([f"- {point}" for point in key_points])
    
    def _format_next_steps(self, next_steps: List[str]) -> str:
        """Format next steps for document inclusion."""
        return '\n'.join([f"- {step}" for step in next_steps]) 